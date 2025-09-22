import sqlite3
from fastapi import APIRouter, HTTPException, UploadFile, File, Depends
from pydantic import BaseModel
import requests
import os
import json
import random
import io
import time
import math
import uuid
from typing import Union, Optional, List, Dict, Any
from bs4 import BeautifulSoup
import re
from docx import Document
from PIL import Image
from pdfminer.high_level import extract_text as pdf_extract_text
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from fastapi.responses import FileResponse
import logging
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from datetime import datetime, timedelta
from blog_routes import router as blog_router
from paddleocr import PaddleOCR
import numpy as np
from urllib.parse import urlparse, parse_qs
from collections import defaultdict
from razorpay_webhook import router as razorpay_router

import os
import tempfile

# Set a writable cache directory for Hugging Face models or other libraries
os.environ["HF_HOME"] = "/tmp/hf_cache"
os.environ["TRANSFORMERS_CACHE"] = "/tmp/hf_cache"
os.environ["XDG_CACHE_HOME"] = "/tmp/cache"

# Ensure the directory exists
os.makedirs("/tmp/hf_cache", exist_ok=True)
os.makedirs("/tmp/cache", exist_ok=True)
# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

router = APIRouter()
security = HTTPBearer(auto_error=False)
from fastapi import FastAPI, Request, HTTPException, APIRouter
from fastapi.responses import JSONResponse
import hmac
import hashlib
import json
import os
from datetime import datetime
from supabase import create_client
from supabase.lib.client_options import ClientOptions

app = FastAPI()
router = APIRouter()

# ---------- Supabase Client ----------
client_options = ClientOptions(
    headers={
        "Authorization": f"Bearer {os.getenv('SUPABASE_SERVICE_ROLE_KEY')}",
        "apikey": os.getenv("SUPABASE_SERVICE_ROLE_KEY"),
        "Content-Type": "application/json",
        "Accept": "application/json"
    }
)

supabase = create_client(
    os.getenv("SUPABASE_URL"),
    os.getenv("SUPABASE_SERVICE_ROLE_KEY"),
    options=client_options
)

# ---------- Razorpay Credentials ----------
RAZORPAY_WEBHOOK_SECRET = os.getenv("RAZORPAY_WEBHOOK_SECRET")

# ---------- Plan to Credit Mapping ----------
CREDIT_LOOKUP = {
    "free": 50,
    "basic": 500,
    "premium": 1000,
    "pro": 999999  # practically unlimited
}

# ---------- Feature Costs ----------
FEATURE_COST = {
    "summary": 5,
    "quiz": 10,
    "humanize": 5,
    "mindmap": 10,
    "diagram": 10,
    "flashcards": 5,
    "vocabulary": 5,
    "handwritten": 5,
    "solve_questions": 5,
    "ppt": 10,
    "study_plan": 5,
    "grammar_correction": 5,
    "knowledge_graph": 10
}

# ---------- Map Razorpay Plan IDs ----------
def map_plan_id_to_internal(plan_id: str):
    mapping = {
        "plan_BASIC_ID": "basic",
        "plan_PREMIUM_ID": "premium",
        "plan_PRO_ID": "pro"
    }
    return mapping.get(plan_id, "free")

# ---------- Verify Razorpay Signature ----------
def verify_signature(payload: bytes, signature: str):
    generated_signature = hmac.new(
        RAZORPAY_WEBHOOK_SECRET.encode(),
        payload,
        hashlib.sha256
    ).hexdigest()
    if generated_signature != signature:
        raise ValueError("Invalid Razorpay webhook signature")

# ---------- Razorpay Webhook ----------
@app.post("/razorpay/webhook")
async def razorpay_webhook(request: Request):
    payload = await request.body()
    signature = request.headers.get("x-razorpay-signature")
    event = json.loads(payload)

    # Verify signature
    try:
        verify_signature(payload, signature)
    except ValueError:
        return JSONResponse({"status": "error", "message": "Invalid signature"}, status_code=400)

    subscription = event.get("payload", {}).get("subscription", {}).get("entity", {})
    user_id = subscription.get("customer_notify_id") or subscription.get("notes", {}).get("user_id")
    plan_id = subscription.get("plan_id")

    internal_plan = map_plan_id_to_internal(plan_id)

    # Update user in Supabase
    supabase.table("users").update({
        "plan": internal_plan,
        "plan_updated_at": datetime.utcnow().isoformat()
    }).eq("id", user_id).execute()

    # Reset credits based on plan
    credits_to_add = CREDIT_LOOKUP.get(internal_plan, 0)
    supabase.table("credits").upsert({
        "user_id": user_id,
        "credits": credits_to_add,
        "last_updated": datetime.utcnow().isoformat()
    }).execute()

    return JSONResponse({"status": "success", "plan": internal_plan, "credits": credits_to_add})

# ---------- Deduct Credit API ----------
@app.post("/deduct-credit")
async def deduct_credit(request: Request):
    body = await request.json()
    user_id = body.get("user_id")
    feature = body.get("feature")

    if not user_id or not feature:
        raise HTTPException(status_code=400, detail="Missing user_id or feature")

    # Fetch credits
    credits_data = supabase.table("credits").select("credits").eq("user_id", user_id).execute()
    if not credits_data.data:
        raise HTTPException(status_code=404, detail="User not found")

    current_credits = credits_data.data[0]["credits"]
    cost = FEATURE_COST.get(feature, 1)

    if current_credits < cost:
        raise HTTPException(status_code=402, detail="Not enough credits")

    new_credits = current_credits - cost

    # Update credits in Supabase
    supabase.table("credits").update({
        "credits": new_credits,
        "last_updated": datetime.utcnow().isoformat()
    }).eq("user_id", user_id).execute()

    return {"success": True, "remaining_credits": new_credits}
router.include_router(blog_router, prefix="/api/blog")
router.include_router(razorpay_router)
# --- Configuration for Groq AI and Fireworks.ai API Calls ---
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
GROQ_API_URL = "https://api.groq.com/openai/v1/chat/completions"
GROQ_MODEL = "openai/gpt-oss-120b"

FIREWORKS_API_KEY = os.getenv("FIREWORKS_API_KEY")
FIREWORKS_API_URL = "https://api.fireworks.ai/inference/v1/chat/completions"
FIREWORKS_MODEL = "accounts/fireworks/models/mixtral-8x7b-instruct"

if not GROQ_API_KEY:
    raise ValueError("GROQ_API_KEY environment variable not set.")

if not FIREWORKS_API_KEY:
    logger.warning("FIREWORKS_API_KEY environment variable not set.")

GROQ_HEADERS = {
    "Authorization": f"Bearer {GROQ_API_KEY}",
    "Content-Type": "application/json",
}

FIREWORKS_HEADERS = {
    "Authorization": f"Bearer {FIREWORKS_API_KEY}",
    "Content-Type": "application/json",
}

YOUTUBE_API_KEY = os.getenv("YOUTUBE_API_KEY")
if not YOUTUBE_API_KEY:
    logger.warning("YOUTUBE_API_KEY environment variable not set.")

MAX_LLM_INPUT_CHARS = 28000

# --- SQLite Database Setup ---
DB_PATH = "/tmp/nexnotes.db"

def init_db():
    with sqlite3.connect(DB_PATH) as conn:
        c = conn.cursor()
        c.execute('''
            CREATE TABLE IF NOT EXISTS user_actions (
                user_id TEXT,
                feature TEXT,
                count INTEGER,
                last_reset TIMESTAMP,
                PRIMARY KEY (user_id, feature)
            )
        ''')
        # c.execute('''
        #     CREATE TABLE IF NOT EXISTS referrals (
        #         user_id TEXT,
        #         feature TEXT,
        #         referral_token TEXT UNIQUE,
        #         used INTEGER DEFAULT 0,
        #         created_at TIMESTAMP
        #     )
        # ''')
        c.execute('''
            CREATE TABLE IF NOT EXISTS guest_usage (
                session_id TEXT PRIMARY KEY,
                usage TEXT,
                updated_at TIMESTAMP
            )
        ''')
        c.execute('INSERT OR IGNORE INTO guest_usage (session_id, usage, updated_at) VALUES (?, ?, ?)',
                  ('guest', json.dumps({}), datetime.now().isoformat()))
        

        c.execute('''
            CREATE TABLE IF NOT EXISTS user_studies (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                user_id TEXT,
                subject TEXT,
                mode TEXT,
                timestamp TIMESTAMP
            )
        ''')
        # Add this to the init_db function in pdf_processor (22).py
        c.execute('''
             CREATE TABLE IF NOT EXISTS user_question_results (
                  id INTEGER PRIMARY KEY AUTOINCREMENT,
                  user_id TEXT,
                  subject TEXT,
                  is_correct BOOLEAN,
                  timestamp TIMESTAMP,
                  FOREIGN KEY (user_id) REFERENCES user_actions (user_id)
            )
        ''')
        
        conn.commit()

init_db()

# --- Pydantic Models ---
class ContentRequest(BaseModel):
    text: str

class GenerateQuestionsRequest(BaseModel):
    text: str
    difficulty: str = "medium"
    count: int = 3

class FollowUpRequest(BaseModel):
    summary: str
    question: str

class URLRequest(BaseModel):
    url: str

class YouTubeURLRequest(BaseModel):
    youtube_url: str

class FlashcardsRequest(BaseModel):
    text: str

class VocabularyRequest(BaseModel):
    text: str

class HumanizeRequest(BaseModel):
    text: str

class MindMapRequest(BaseModel):
    text: str

class DiagramRequest(BaseModel):
    text: str
    diagram_type: str = "flowchart"

class HandwrittenRequest(BaseModel):
    text: str
    style: str = "neat"

class SolveQuestionsRequest(BaseModel):
    text: str

class AnswerSubmissionRequest(BaseModel):
    user_id: Optional[str]
    subject: str
    answers: List[dict]  # List of { question_idx: int, user_answer: str, correct_answer: str }

class ContentRequest(BaseModel):
    text: str
    tone: str = "medium"  # Add tone parameter with default value
class PPTRequest(BaseModel):
    text: str
    slide_count: int = 3
    font_family: str = "Arial"
    font_color: str = "#000000"
    slide_bg_color: str = "#FFFFFF"
    font_size: int = 18

class StudyPlanRequest(BaseModel):
    text: str
    duration_days: int = 7
    hours_per_day: int = 2

class GrammarCorrectionRequest(BaseModel):
    text: str

# class ReferralRequest(BaseModel):
#     feature: str
#     user_id: str

# class ReferralValidationRequest(BaseModel):
#     referral_token: str

# class ReferralStatusRequest(BaseModel):
#     user_id: str
#     feature: str

class GuestUsageRequest(BaseModel):
    usage: dict

class BlogPostCreate(BaseModel):
    title: str
    content: str
    author: str = "Admin"
    tags: List[str] = []

class BlogPostResponse(BaseModel):
    id: str
    title: str
    content: str
    author: str
    tags: List[str]
    created_at: datetime
    updated_at: datetime

# --- Helper Functions ---
async def get_user_id(credentials: HTTPAuthorizationCredentials = Depends(security)) -> Optional[str]:
    if credentials is None:
        return None  # Guest user
    return credentials.credentials

def call_llm(prompt_messages: list, temperature: float = 0.7, max_tokens: int = 1024, retries: int = 5) -> str:
    groq_payload = {
        "model": GROQ_MODEL,
        "messages": prompt_messages,
        "temperature": temperature,
        "max_tokens": max_tokens,
        "top_p": 0.9
    }

    fireworks_payload = {
        "model": FIREWORKS_MODEL,
        "messages": prompt_messages,
        "temperature": temperature,
        "max_tokens": max_tokens,
        "top_p": 0.9,
        "repetition_penalty": 1.0,
        "top_k": 50
    }

    for i in range(retries):
        try:
            response = requests.post(GROQ_API_URL, headers=GROQ_HEADERS, json=groq_payload, timeout=30)
            response.raise_for_status()
            result = response.json()
            if not result or "choices" not in result or not result["choices"]:
                raise Exception("Invalid response from Groq AI: No choices found.")
            return result["choices"][0]["message"]["content"]
        except requests.exceptions.HTTPError as e:
            if e.response.status_code == 429 and FIREWORKS_API_KEY:
                logger.warning(f"Groq AI rate limit hit (429). Falling back to Fireworks.ai... (Attempt {i + 1}/{retries})")
                try:
                    response = requests.post(FIREWORKS_API_URL, headers=FIREWORKS_HEADERS, json=fireworks_payload, timeout=30)
                    response.raise_for_status()
                    result = response.json()
                    if not result or "choices" not in result or not result["choices"]:
                        raise Exception("Invalid response from Fireworks.ai: No choices found.")
                    return result["choices"][0]["message"]["content"]
                except requests.exceptions.HTTPError as fw_e:
                    if fw_e.response.status_code == 429:
                        wait_time = min(2 ** i, 60)
                        logger.warning(f"Fireworks.ai rate limit hit (429). Retrying in {wait_time} seconds...")
                        time.sleep(wait_time)
                    elif fw_e.response.status_code == 401:
                        raise HTTPException(status_code=401, detail="Unauthorized: Invalid Fireworks.ai API key provided.")
                    else:
                        logger.error(f"Error from Fireworks.ai: {fw_e.response.status_code} {fw_e.response.reason}")
                        raise HTTPException(status_code=500, detail=f"Unexpected error from Fireworks.ai: {fw_e.response.status_code}")
                except Exception as fw_e:
                    logger.error(f"Unexpected error from Fireworks.ai: {fw_e}")
                    raise HTTPException(status_code=500, detail=f"Failed to connect to Fireworks.ai API: {fw_e}")
            elif e.response.status_code == 429:
                wait_time = min(2 ** i, 60)
                logger.warning(f"Groq AI rate limit hit (429), no Fireworks.ai key. Retrying in {wait_time} seconds...")
                time.sleep(wait_time)
            elif e.response.status_code == 400:
                error_response = e.response.json()
                logger.error(f"Client error from Groq AI: {e.response.status_code} {e.response.reason}")
                raise HTTPException(status_code=400, detail=f"Bad Request: {error_response.get('error', {}).get('message', 'Invalid request')}")
            elif e.response.status_code == 401:
                raise HTTPException(status_code=401, detail="Unauthorized: Invalid Groq AI API key provided.")
            else:
                logger.error(f"Server error from Groq AI: {e.response.status_code} {e.response.reason}")
                raise HTTPException(status_code=500, detail=f"Unexpected error from Groq AI: {e.response.status_code}")
        except requests.exceptions.RequestException as e:
            logger.error(f"Network error making API call to Groq AI: {e}")
            raise HTTPException(status_code=500, detail=f"Failed to connect to Groq AI API: {e}")
        except Exception as e:
            logger.error(f"An unexpected error occurred during Groq AI API call: {e}")
            raise HTTPException(status_code=500, detail=f"Internal Server Error during AI call: {e}")
    raise HTTPException(status_code=500, detail="Failed to get a response from Groq AI or Fireworks.ai after multiple retries.")

# --- Helper Functions for Text Extraction ---
async def extract_text_from_pdf(pdf_file: UploadFile) -> str:
    try:
        pdf_bytes = await pdf_file.read()
        text = pdf_extract_text(io.BytesIO(pdf_bytes))
        if not text.strip():
            logger.warning(f"No text extracted from PDF '{pdf_file.filename}'")
            return ""
        return text
    except Exception as e:
        logger.error(f"Error extracting text from PDF '{pdf_file.filename}': {e}")
        raise HTTPException(status_code=500, detail=f"Failed to extract text from PDF: {e}")

async def extract_text_from_docx(docx_file: UploadFile) -> str:
    try:
        docx_bytes = await docx_file.read()
        document = Document(io.BytesIO(docx_bytes))
        full_text = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        extracted_content = '\n'.join(full_text)
        if not extracted_content.strip():
            logger.warning(f"No text extracted from DOCX '{docx_file.filename}'")
            return ""
        return extracted_content
    except Exception as e:
        logger.error(f"Error extracting text from DOCX '{docx_file.filename}': {e}")
        raise HTTPException(status_code=500, detail=f"Failed to extract text from Word document: {e}")

async def extract_text_from_ppt(ppt_file: UploadFile) -> str:
    try:
        ppt_bytes = await ppt_file.read()
        prs = Presentation(io.BytesIO(ppt_bytes))
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    full_text.append(shape.text)
        extracted_content = '\n'.join(full_text)
        if not extracted_content.strip():
            logger.warning(f"No text extracted from PPT '{ppt_file.filename}'")
            return ""
        return extracted_content
    except Exception as e:
        logger.error(f"Error extracting text from PPT '{ppt_file.filename}': {e}")
        raise HTTPException(status_code=500, detail=f"Failed to extract text from PowerPoint document: {e}")

async def extract_text_from_image(image_file: UploadFile) -> str:
    try:
        model_dir = "/tmp/paddleocr"
        os.makedirs(model_dir, exist_ok=True)
        ocr = PaddleOCR(
            use_angle_cls=True,
            lang='en',
            use_gpu=False,
            det_model_dir=f"{model_dir}/det/en/en_PP-OCRv3_det_infer",
            rec_model_dir=f"{model_dir}/rec/en/en_PP-OCRv3_rec_infer",
            cls_model_dir=f"{model_dir}/cls/cls_mv3_infer"
        )
        image_bytes = await image_file.read()
        image = Image.open(io.BytesIO(image_bytes)).convert("RGB")
        image_np = np.array(image)
        result = ocr.ocr(image_np, cls=True)
        if not result or not result[0]:
            logger.warning(f"No text extracted from image '{image_file.filename}'")
            return ""
        extracted_text = ' '.join([line[1][0] for line in result[0] if line[1][0]]).strip()
        if not extracted_text:
            logger.warning(f"No valid text extracted from image '{image_file.filename}'")
            return ""
        if extracted_text.isdigit() or extracted_text == "0 0":
            logger.warning(f"Invalid text extracted from image '{image_file.filename}': '{extracted_text}'")
            return ""
        return extracted_text
    except Exception as e:
        logger.error(f"Error extracting text from image '{image_file.filename}': {e}")
        raise HTTPException(status_code=500, detail=f"Failed to extract text from image: {e}")



def extract_youtube_transcript(youtube_url: str) -> str:
    try:
        video_id = parse_qs(urlparse(youtube_url).query).get('v', [None])[0]
        if not video_id:
            raise ValueError("Invalid YouTube URL provided.")
        api_url = f"https://www.googleapis.com/youtube/v3/captions?part=snippet&videoId={video_id}&key={YOUTUBE_API_KEY}"
        response = requests.get(api_url)
        response.raise_for_status()
        captions_data = response.json()
        if not captions_data.get("items"):
            raise ValueError("No captions found for this video.")
        caption_id = captions_data["items"][0]["id"]
        caption_download_url = f"https://www.googleapis.com/youtube/v3/captions/{caption_id}?tfmt=srt&key={YOUTUBE_API_KEY}"
        caption_response = requests.get(caption_download_url)
        caption_response.raise_for_status()
        transcript = caption_response.text
        if not transcript.strip():
            logger.warning(f"No transcript extracted from YouTube URL '{youtube_url}'")
            return ""
        return transcript
    except Exception as e:
        logger.error(f"Error extracting transcript from YouTube URL '{youtube_url}': {e}")
        raise HTTPException(status_code=500, detail=f"Failed to extract transcript from YouTube: {e}")

def generate_questions(text: str, difficulty: str, count: int) -> List[dict]:
    prompt_content = f"""Generate {count} {difficulty} difficulty questions from the following content.
Each question should be a multiple-choice question with 4 options (A, B, C, D), one correct answer, and an explanation.
Format the output strictly as a JSON object with a single key "questions", whose value is a list of question objects.
Each question object should have 'question', 'options' (a list of 4 strings), 'correct_answer' (a single letter A/B/C/D), and 'explanation'.
Do not include any other text or markdown outside the JSON.

Content:
{text.strip()}
"""
    prompt_messages = [
        {"role": "system", "content": "You are an educational assistant specialized in generating multiple-choice questions in strict JSON format."},
        {"role": "user", "content": prompt_content}
    ]

    try:
        raw_response = call_llm(prompt_messages, temperature=0.7, max_tokens=1500)
        if raw_response.startswith("```json") and raw_response.endswith("```"):
            raw_response = raw_response[len("```json"):-len("```")].strip()
        elif raw_response.startswith("```") and raw_response.endswith("```"):
            raw_response = raw_response[len("```"):-len("```")].strip()
        
        questions_data = json.loads(raw_response)
        if not isinstance(questions_data, dict) or "questions" not in questions_data or not isinstance(questions_data["questions"], list):
            raise ValueError("Model did not return a valid JSON object with a 'questions' list.")
        
        final_questions = []
        for q in questions_data["questions"]:
            if not all(k in q for k in ["question", "options", "correct_answer", "explanation"]):
                continue
            if not isinstance(q["options"], list) or len(q["options"]) != 4:
                continue
            if q["correct_answer"] not in ['A', 'B', 'C', 'D']:
                continue
            final_questions.append({
                "question": q["question"].strip(),
                "options": [opt.strip() for opt in q["options"]],
                "correct_answer": q["correct_answer"],
                "explanation": q["explanation"].strip()
            })
        
        if len(final_questions) != count:
            raise HTTPException(status_code=400, detail=f"Expected {count} questions, but generated {len(final_questions)} valid questions.")
        
        return final_questions
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=500, detail=f"Failed to parse JSON from AI model: {e}")
    except ValueError as e:
        raise HTTPException(status_code=500, detail=f"Invalid data format from AI model: {e}")
    except HTTPException as e:
        raise e
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate questions: {str(e)}")

def generate_flashcards(text: str) -> List[dict]:
    prompt_content = f"""Generate 5-10 flashcards from the following content.
Each flashcard should have a 'front' (question or term) and 'back' (answer or definition).
Format the output strictly as a JSON object with a single key "flashcards", whose value is a list of flashcard objects.
Each flashcard object should have 'front' and 'back'.
Do not include any other text or markdown outside the JSON.

Content:
{text.strip()}
"""
    prompt_messages = [
        {"role": "system", "content": "You are an educational assistant specialized in generating flashcards in strict JSON format."},
        {"role": "user", "content": prompt_content}
    ]

    try:
        raw_response = call_llm(prompt_messages, temperature=0.6, max_tokens=1000)
        if raw_response.startswith("```json") and raw_response.endswith("```"):
            raw_response = raw_response[len("```json"):-len("```")].strip()
        elif raw_response.startswith("```") and raw_response.endswith("```"):
            raw_response = raw_response[len("```"):-len("```")].strip()
        
        flashcards_data = json.loads(raw_response)
        if not isinstance(flashcards_data, dict) or "flashcards" not in flashcards_data or not isinstance(flashcards_data["flashcards"], list):
            raise ValueError("Model did not return a valid JSON object with a 'flashcards' list.")
        
        final_flashcards = []
        for fc in flashcards_data["flashcards"]:
            if not all(k in fc for k in ["front", "back"]):
                continue
            final_flashcards.append({
                "front": fc["front"].strip(),
                "back": fc["back"].strip()
            })
        
        if not final_flashcards:
            raise HTTPException(status_code=400, detail="No valid flashcards could be generated from the provided content.")
        
        return final_flashcards
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=500, detail=f"Failed to parse JSON from AI model: {e}")
    except ValueError as e:
        raise HTTPException(status_code=500, detail=f"Invalid data format from AI model: {e}")
    except HTTPException as e:
        raise e
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate flashcards: {str(e)}")

def generate_vocabulary(text: str) -> List[dict]:
    prompt_content = f"""Extract 10-15 key vocabulary words from the following content.
For each word, provide a definition and an example sentence from or related to the content.
Format the output strictly as a JSON object with a single key "vocabulary", whose value is a list of vocabulary objects.
Each vocabulary object should have 'word', 'definition', and 'example'.
Do not include any other text or markdown outside the JSON.

Content:
{text.strip()}
"""
    prompt_messages = [
        {"role": "system", "content": "You are an educational assistant specialized in extracting vocabulary in strict JSON format."},
        {"role": "user", "content": prompt_content}
    ]

    try:
        raw_response = call_llm(prompt_messages, temperature=0.5, max_tokens=1000)
        if raw_response.startswith("```json") and raw_response.endswith("```"):
            raw_response = raw_response[len("```json"):-len("```")].strip()
        elif raw_response.startswith("```") and raw_response.endswith("```"):
            raw_response = raw_response[len("```"):-len("```")].strip()
        
        vocab_data = json.loads(raw_response)
        if not isinstance(vocab_data, dict) or "vocabulary" not in vocab_data or not isinstance(vocab_data["vocabulary"], list):
            raise ValueError("Model did not return a valid JSON object with a 'vocabulary' list.")
        
        final_vocab = []
        for v in vocab_data["vocabulary"]:
            if not all(k in v for k in ["word", "definition", "example"]):
                continue
            final_vocab.append({
                "word": v["word"].strip(),
                "definition": v["definition"].strip(),
                "example": v["example"].strip()
            })
        
        if not final_vocab:
            raise HTTPException(status_code=400, detail="No valid vocabulary could be extracted from the provided content.")
        
        return final_vocab
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=500, detail=f"Failed to parse JSON from AI model: {e}")
    except ValueError as e:
        raise HTTPException(status_code=500, detail=f"Invalid data format from AI model: {e}")
    except HTTPException as e:
        raise e
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate vocabulary: {str(e)}")

def humanize_text(text: str) -> str:
    prompt_messages = [
        {"role": "system", "content": "You are a friendly educational assistant that rewrites content in a natural, human-like way. Use Markdown for formatting."},
        {"role": "user", "content": f"""You are an ultra professional human editor and writer. Rewrite the given text in a very natural, human-like way while preserving its meaning. 
        -write it like it has been written in simplest words possible , very basic and very natural and someone with average grammar
        - dont repeat words often, dont use modern words.
        - use very very basic english, no modernity.
        - dont be too generic
        - dont use the words commonly used by AI
         - WRITE LIKE YOU ARE NEW TO WRITING ENGLISH, DONT WRITE IT PERFECTLY AND IN VERY SYSTEMATIC GRAMMAR
         - ADD WORDS LIKE honestly, to be fair, lets face it, etc.
         - do grammatical imperfections, use as simplest terms as possible.DO NOT USE LARGE OR DIFFICULT WORDS
         - Vary sentence length and structure.
         - Use everyday expressions and natural transitions.
         - Avoid repetitive or overly formal phrasing.
         - Include a few subtle human touches like personal tone, rhetorical questions, or conversational flow. keep the text as simple as possible
         - Make it indistinguishable from a text written manually by a skilled human.
         SEND THE HUMANIZED TEXT ONLY, NOTHING ELSE



Content to humanize:
{text.strip()}
"""}
    ]
    
    try:
        humanized = call_llm(prompt_messages, temperature=0.8, max_tokens=1024)
        return humanized.strip()
    except HTTPException as e:
        raise e
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to humanize text: {str(e)}")

def correct_grammar(text: str) -> str:
    prompt_messages = [
        {"role": "system", "content": "You are a grammar correction assistant. Correct grammar, spelling, and punctuation while preserving the original meaning. Use Markdown for formatting if needed."},
        {"role": "user", "content": f"""Correct the grammar, spelling, and punctuation in the following text.
Do not change the meaning or add/remove information. Return only the corrected text.

Text:
{text.strip()}
"""}
    ]
    
    try:
        corrected = call_llm(prompt_messages, temperature=0.3, max_tokens=1024)
        return corrected.strip()
    except HTTPException as e:
        raise e
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to correct grammar: {str(e)}")


@router.post("/api/upload-and-extract")
async def upload_and_extract(file: UploadFile = File(...), user_id: Optional[str] = Depends(get_user_id)):
   
    filename = file.filename.lower()
    extracted_text = ""

    if filename.endswith('.pdf'):
        extracted_text = await extract_text_from_pdf(file)
    elif filename.endswith('.docx'):
        extracted_text = await extract_text_from_docx(file)
    elif filename.endswith('.doc'):
        raise HTTPException(status_code=400, detail="Unsupported file type: .doc")
    elif filename.endswith(('.ppt', '.pptx')):
        extracted_text = await extract_text_from_ppt(file)
    elif filename.endswith(('.jpg', '.jpeg', '.png')):
        extracted_text = await extract_text_from_image(file)
    else:
        raise HTTPException(status_code=400, detail=f"Unsupported file type: {file.filename}")

    if not extracted_text.strip():
        raise HTTPException(status_code=400, detail="No readable text found in the uploaded file.")

    if len(extracted_text) > MAX_LLM_INPUT_CHARS:
        extracted_text = extracted_text[:MAX_LLM_INPUT_CHARS]
        logger.warning(f"Text truncated to {MAX_LLM_INPUT_CHARS} characters.")

    return {"extracted_text": extracted_text}

@router.post("/api/fetch-and-extract-url")
async def fetch_and_extract_url(request: URLRequest, user_id: Optional[str] = Depends(get_user_id)):
    
    url = request.url.strip()
    if not url:
        raise HTTPException(status_code=400, detail="No URL provided.")

    if not (url.startswith('http://') or url.startswith('https://')):
        url = 'https://' + url

    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        soup = BeautifulSoup(response.text, 'lxml')
        for element in soup(["script", "style", "nav", "footer", "header", "form", "aside", "meta", "iframe", "img", "svg", "link", "input", "button", "select", "textarea"]):
            element.extract()

        MAIN_CONTENT_SELECTORS = [
            'article', 'main', 'div.main-content', 'div.entry-content', 'div.post-content',
            'div.article-body', 'div[role="main"]', 'div#content', 'div#main', 'div.content',
            'div.post', 'div.blog-post', 'div.article', 'p'
        ]

        main_text = ""
        best_candidate_text = ""
        for selector in MAIN_CONTENT_SELECTORS:
            elements = soup.select(selector)
            for element in elements:
                current_text = element.get_text(separator=' ', strip=True)
                current_text = re.sub(r'\s+', ' ', current_text).strip()
                current_text = re.sub(r'[\xa0\u200b]+', ' ', current_text).strip()
                if len(current_text) > len(best_candidate_text):
                    best_candidate_text = current_text

        if not best_candidate_text or len(best_candidate_text) < 100:
            body_element = soup.find('body')
            if body_element:
                for element in body_element(["script", "style", "nav", "footer", "header", "form", "aside", "meta", "iframe", "img", "svg", "link", "input", "button", "select", "textarea"]):
                    element.extract()
                main_text = body_element.get_text(separator=' ', strip=True)
                main_text = re.sub(r'\s+', ' ', main_text).strip()
                main_text = re.sub(r'[\xa0\u200b]+', ' ', main_text).strip()
            else:
                main_text = soup.get_text(separator=' ', strip=True)
                main_text = re.sub(r'\s+', ' ', main_text).strip()
                main_text = re.sub(r'[\xa0\u200b]+', ' ', main_text).strip()
        else:
            main_text = best_candidate_text

        if not main_text.strip():
            raise HTTPException(status_code=400, detail="Could not extract any meaningful text from the URL.")

        if len(main_text) > MAX_LLM_INPUT_CHARS:
            main_text = main_text[:MAX_LLM_INPUT_CHARS]
            logger.warning(f"Text truncated to {MAX_LLM_INPUT_CHARS} characters.")

        return {"extracted_text": main_text}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to extract text from URL: {e}")

import re

def sanitize_mermaid_mindmap(code: str) -> str:
    # Remove Markdown code fences and language labels
    code = re.sub(r"```[a-zA-Z0-9_-]*", "", code)
    code = code.replace("```", "")
    code = code.strip()

    # Ensure "mindmap" is alone on the first line
    lines = code.splitlines()
    cleaned_lines = []
    found_header = False
    for line in lines:
        stripped = line.strip()
        if stripped.startswith("mindmap") and not found_header:
            cleaned_lines.append("mindmap")
            found_header = True
        elif stripped.startswith("mindmap") and found_header:
            # skip duplicate "mindmap" lines
            continue
        else:
            cleaned_lines.append(stripped)
    code = "\n".join(cleaned_lines)

    # Ensure root node is on its own line
    code = re.sub(r"^mindmap\s+root", "mindmap\n  root", code, flags=re.MULTILINE)

    # Fix missing parentheses in root node
    lines = code.splitlines()
    for i, line in enumerate(lines):
        if line.strip().startswith("root") and not re.search(r"\(\(.*\)\)", line):
            content = re.sub(r"^root\s*\(*\)*", "", line).strip() or "Main Topic"
            lines[i] = f"  root(({content}))"
    return "\n".join(lines).strip()

# def fix_mindmap_indentation(code: str) -> str:
#     """
#     Auto-indent mermaid mindmap code so that branches, leaves, and subleaves
#     render properly.
#     """
#     lines = code.split("\n")
#     fixed = []
#     for line in lines:
#         clean = line.strip()
#         if not clean:
#             continue
#         if clean.startswith("root"):
#             fixed.append("  " + clean)
#         elif clean.lower().startswith("branch"):
#             fixed.append("    " + clean)
#         elif clean.lower().startswith("leaf"):
#             fixed.append("      " + clean)
#         elif clean.lower().startswith("subleaf"):
#             fixed.append("        " + clean)
#         else:
#             fixed.append(clean)
#     return "\n".join(["mindmap"] + fixed)

@router.post("/api/generate-mindmap")
async def generate_mindmap(request: MindMapRequest, user_id: Optional[str] = Depends(get_user_id)):
    if not request.text.strip():
        raise HTTPException(status_code=400, detail="No text provided for mind map generation.")

    prompt_messages = [
        {"role": "system", "content": """You are an educational assistant that generates mind maps in strict Mermaid.js format. 
Your output MUST follow this exact format:
mindmap
  root((Main Topic))
    Branch A
      Leaf 1
      Leaf 2


Only output the Mermaid code, nothing else.Make sure it covers all the major topics of the provided content and its context. DO NOT include ```mermaid or extra text."""},
        {"role": "user", "content": f"Generate a comprehensive mind map in Mermaid.js format from:\n{request.text.strip()}"}
    ]
    
    try:
        mermaid_code = call_llm(prompt_messages, temperature=0.3, max_tokens=1024)
        mermaid_code = sanitize_mermaid_mindmap(mermaid_code)
        return {"mermaid_code": mermaid_code}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate mind map: {str(e)}")


@router.post("/api/generate-diagram")
async def generate_diagram(request: DiagramRequest, user_id: Optional[str] = Depends(get_user_id)):
    if not request.text.strip():
        raise HTTPException(status_code=400, detail="No text provided for diagram generation.")

    diagram_types = {
        "flowchart": "flowchart TD",
        "sequence": "sequenceDiagram",
        "class": "classDiagram",
        "state": "stateDiagram-v2",
        "entity": "erDiagram"
    }

    diagram_type = diagram_types.get(request.diagram_type, "flowchart TD")
    diagram_type_key = request.diagram_type if request.diagram_type in diagram_types else "flowchart"

    prompt_content = f"""You are an expert in generating {diagram_type_key} diagrams in strict Mermaid.js format. Your task is to create a valid {diagram_type_key} diagram based on the provided content. Follow these rules:
1. Output ONLY valid Mermaid.js code starting with '{diagram_type}'.
2. Ensure the diagram is syntactically correct and can be parsed by Mermaid.js.
3. Use concise labels and avoid special characters like quotes, brackets, or colons in node names unless required by the diagram type.
4. For short or ambiguous input (e.g., a single word), create a simple valid diagram with the input as a central node or actor and generic connections (e.g., 'Node1', 'Node2' for flowcharts, or 'Actor1', 'Actor2' for sequence diagrams).
5. For flowcharts, use 'id[Label]' for nodes and 'id1-->id2' for connections.
6. For sequence diagrams, use 'actor ActorName' and 'Actor1->>Actor2: Message'.
7. For class diagrams, use 'class ClassName {{}}' and 'Class1--Class2' for relationships.
8. For state diagrams, use '[*] --> State1' and 'State1 --> State2'.
9. For entity-relationship diagrams, use 'ENTITY ||--o{{ RELATION : "description"'.
10. Do NOT include explanations, markdown, or code fences (```).

Example for flowchart:
flowchart TD
  A[Main Topic] --> B[Concept 1]
  A --> C[Concept 2]
  B --> D[Detail 1]

Example for sequence:
sequenceDiagram
  actor User
  actor System
  User->>System: Request
  System-->>User: Response

Content to generate {diagram_type_key} diagram from:
{request.text.strip()}"""

    prompt_messages = [
        {"role": "system", "content": f"You are an educational assistant that generates {diagram_type_key} diagrams in strict Mermaid.js format. Output ONLY the Mermaid code without any additional text or markdown. DO NOT INCLUDE SPECIAL CHARACTERS IN CONTENTS UNLESS REQUIRED."},
        {"role": "user", "content": prompt_content}
    ]
    
    try:
        mermaid_code = call_llm(prompt_messages, temperature=0.3, max_tokens=1024)
        if not mermaid_code.strip().startswith(diagram_type.split()[0]):
            fallback_code = f"{diagram_type}\n"
            if diagram_type_key == "flowchart":
                fallback_code += f"  A[{request.text.strip()[:50] or 'Main Topic'}] --> B[Concept 1]\n  A --> C[Concept 2]"
            elif diagram_type_key == "sequence":
                fallback_code += f"  actor A\n  actor B\n  A->>B: Message"
            elif diagram_type_key == "class":
                fallback_code += f"  class {request.text.strip()[:50] or 'MainClass'} {{}}\n  class RelatedClass {{}}\n  MainClass--RelatedClass"
            elif diagram_type_key == "state":
                fallback_code += f"  [*] --> {request.text.strip()[:50] or 'State1'}\n  {request.text.strip()[:50] or 'State1'} --> State2"
            elif diagram_type_key == "entity":
                fallback_code += f"  {request.text.strip().upper()[:50] or 'ENTITY1'} ||--o{{ ENTITY2 : relates"
            mermaid_code = fallback_code
        lines = mermaid_code.strip().splitlines()
        cleaned_lines = [lines[0]]
        for line in lines[1:]:
            cleaned_line = re.sub(r'[^a-zA-Z0-9\s\-\[\]\(\)\{\}\:\>\<\|\*]', '', line.strip())
            if cleaned_line:
                cleaned_lines.append(cleaned_line)
        mermaid_code = "\n".join(cleaned_lines)
        return {"mermaid_code": mermaid_code.strip()}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate diagram: {str(e)}")


   



def generate_handwritten(text: str, style: str) -> str:
    valid_styles = ['neat', 'casual', 'messy']
    if style not in valid_styles:
        raise HTTPException(status_code=400, detail=f"Invalid handwritten style. Supported: {', '.join(valid_styles)}")
    
    prompt_messages = [
        {"role": "system", "content": f"You are an assistant that simulates {style} handwritten notes. Format the output as plain text with line breaks to mimic handwriting."},
        {"role": "user", "content": f"""Rewrite the following content as {style} handwritten notes.
Use line breaks and simple formatting to simulate handwriting. Do not use Markdown. ONLY RETURN THE CONTENT PROVIDED BY USER. DO NOT ADD OR REMOVE ANYTHING ELSE.

Content:
{text.strip()}
"""}
    ]
    
    try:
        handwritten = call_llm(prompt_messages, temperature=0.9 if style == 'messy' else 0.7, max_tokens=1024)
        return handwritten.strip()
    except HTTPException as e:
        raise e
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate handwritten notes: {str(e)}")

def solve_questions(text: str) -> List[dict]:
    prompt_content = f"""Solve the questions in the following content.
For each question, provide a clear answer and explanation.
Format the output strictly as a JSON object with a single key "answers", whose value is a list of answer objects.
Each answer object should have 'question' and 'answer'.
Do not include any other text or markdown outside the JSON.

Content with questions:
{text.strip()}
"""
    prompt_messages = [
        {"role": "system", "content": "You are an educational assistant specialized in solving questions in strict JSON format."},
        {"role": "user", "content": prompt_content}
    ]

    try:
        raw_response = call_llm(prompt_messages, temperature=0.4, max_tokens=1500)
        if raw_response.startswith("```json") and raw_response.endswith("```"):
            raw_response = raw_response[len("```json"):-len("```")].strip()
        elif raw_response.startswith("```") and raw_response.endswith("```"):
            raw_response = raw_response[len("```"):-len("```")].strip()
        
        answers_data = json.loads(raw_response)
        if not isinstance(answers_data, dict) or "answers" not in answers_data or not isinstance(answers_data["answers"], list):
            raise ValueError("Model did not return a valid JSON object with an 'answers' list.")
        
        final_answers = []
        for a in answers_data["answers"]:
            if not all(k in a for k in ["question", "answer"]):
                continue
            final_answers.append({
                "question": a["question"].strip(),
                "answer": a["answer"].strip()
            })
        
        if not final_answers:
            raise HTTPException(status_code=400, detail="No valid answers could be generated from the provided content.")
        
        return final_answers
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=500, detail=f"Failed to parse JSON from AI model: {e}")
    except ValueError as e:
        raise HTTPException(status_code=500, detail=f"Invalid data format from AI model: {e}")
    except HTTPException as e:
        raise e
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to solve questions: {str(e)}")

def generate_ppt(text: str, slide_count: int, font_family: str, font_color: str, slide_bg_color: str, font_size: int) -> str:
    if slide_count < 1 or slide_count > 20:
        raise HTTPException(status_code=400, detail="Slide count must be between 1 and 20.")
    
    prompt_content = f"""Generate content for a PowerPoint presentation with {slide_count} slides based on the following text.
For each slide, provide 'title' and 'content' (a list of bullet points).
Format the output strictly as a JSON object with a single key "slides", whose value is a list of slide objects.
Each slide object should have 'title' and 'content' (list of strings).
Do not include any other text or markdown outside the JSON.

Content:
{text.strip()}
"""
    prompt_messages = [
        {"role": "system", "content": "You are an assistant specialized in generating PPT slide content in strict JSON format."},
        {"role": "user", "content": prompt_content}
    ]

    try:
        raw_response = call_llm(prompt_messages, temperature=0.6, max_tokens=1500)
        if raw_response.startswith("```json") and raw_response.endswith("```"):
            raw_response = raw_response[len("```json"):-len("```")].strip()
        elif raw_response.startswith("```") and raw_response.endswith("```"):
            raw_response = raw_response[len("```"):-len("```")].strip()
        
        slides_data = json.loads(raw_response)
        if not isinstance(slides_data, dict) or "slides" not in slides_data or not isinstance(slides_data["slides"], list):
            raise ValueError("Model did not return a valid JSON object with a 'slides' list.")
        
        final_slides = []
        for slide in slides_data["slides"]:
            if not all(k in slide for k in ["title", "content"]):
                continue
            if not isinstance(slide["content"], list):
                continue
            final_slides.append({
                "title": slide["title"].strip(),
                "content": [point.strip() for point in slide["content"]]
            })
        
        if len(final_slides) != slide_count:
            raise HTTPException(status_code=400, detail=f"Expected {slide_count} slides, but generated {len(final_slides)} valid slides.")
        
        # Generate PPT file
        prs = Presentation()
        r, g, b = int(font_color[1:3], 16), int(font_color[3:5], 16), int(font_color[5:7], 16)
        bg_r, bg_g, bg_b = int(slide_bg_color[1:3], 16), int(slide_bg_color[3:5], 16), int(slide_bg_color[5:7], 16)
        
        for slide_data in final_slides:
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = slide_data["title"]
            title.text_frame.paragraphs[0].font.size = Pt(font_size + 6)
            title.text_frame.paragraphs[0].font.name = font_family
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(r, g, b)
            
            content = slide.placeholders[1]
            for point in slide_data["content"]:
                p = content.text_frame.add_paragraph()
                p.text = point
                p.font.size = Pt(font_size)
                p.font.name = font_family
                p.font.color.rgb = RGBColor(r, g, b)
            
            # Set background color
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(bg_r, bg_g, bg_b)
        
        ppt_path = f"/tmp/ppt-{uuid.uuid4()}.pptx"
        prs.save(ppt_path)
        return ppt_path
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=500, detail=f"Failed to parse JSON from AI model: {e}")
    except ValueError as e:
        raise HTTPException(status_code=500, detail=f"Invalid data format from AI model: {e}")
    except HTTPException as e:
        raise e
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate PPT: {str(e)}")

      
    

def generate_study_plan(request: StudyPlanRequest) -> dict:
    prompt_content = f"""Generate a study plan for {request.duration_days} days, with {request.hours_per_day} daily study hours. Follow these rules:
1. Break down the content into key topics or sections. Tell what things to study or revise each day
2. Distribute the topics across the specified number of days ({request.duration_days}), ensuring an even workload.
3. For each day, specify the study duration (within {request.hours_per_day} hours total) and include specific tasks (e.g., reading, summarizing, practicing questions, reviewing flashcards).
4. Include review sessions and practice tasks to reinforce learning.
5. Format the output strictly as a JSON object with a single key "study_plan", whose value is a list of daily plans. Each daily plan should have 'day', 'duration', and 'tasks' (a list of task objects with 'task' and 'description').
6. Do not include any other text or markdown outside the JSON.

Example JSON format:
{{
  "study_plan": [
    {{
      "day": 1,
      "duration": "2 hours",
      "tasks": [
        {{
          "task": "Read Section 1",
          "description": "Read and take notes on the introduction to photosynthesis."
        }},
        {{
          "task": "Practice Questions",
          "description": "Answer 5 multiple-choice questions on photosynthesis."
        }}
      ]
    }},
    {{
      "day": 2,
      "duration": "2 hours",
      "tasks": [
        {{
          "task": "Review Section 1",
          "description": "Review notes on photosynthesis and clarify doubts."
        }},
        {{
          "task": "Read Section 2",
          "description": "Read about the Calvin cycle."
        }}
      ]
    }}
  ]
}}

Content to generate study plan from:
{request.text.strip()}
"""
    prompt_messages = [
        {"role": "system", "content": "You are an educational assistant specialized in generating systematic study plans in strict JSON format."},
        {"role": "user", "content": prompt_content}
    ]

    try:
        raw_response = call_llm(prompt_messages, temperature=0.5, max_tokens=1500)
        if raw_response.startswith("```json") and raw_response.endswith("```"):
            raw_response = raw_response[len("```json"):-len("```")].strip()
        elif raw_response.startswith("```") and raw_response.endswith("```"):
            raw_response = raw_response[len("```"):-len("```")].strip()
        
        study_plan_data = json.loads(raw_response)
        if not isinstance(study_plan_data, dict) or "study_plan" not in study_plan_data or not isinstance(study_plan_data["study_plan"], list):
            raise ValueError("Model did not return a valid JSON object with a 'study_plan' list.")
        
        final_study_plan = []
        for day in study_plan_data["study_plan"]:
            if not all(k in day for k in ["day", "duration", "tasks"]):
                continue
            if not isinstance(day["tasks"], list):
                continue
            valid_tasks = []
            for task in day["tasks"]:
                if not all(k in task for k in ["task", "description"]):
                    continue
                valid_tasks.append({
                    "task": task["task"].strip(),
                    "description": task["description"].strip()
                })
            if valid_tasks:
                final_study_plan.append({
                    "day": day["day"],
                    "duration": day["duration"].strip(),
                    "tasks": valid_tasks
                })
        
        if not final_study_plan:
            raise HTTPException(status_code=400, detail="No valid study plan could be generated from the provided content.")
        
        return {"study_plan": final_study_plan}
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=500, detail=f"Failed to parse JSON from AI model: {e}")
    except ValueError as e:
        raise HTTPException(status_code=500, detail=f"Invalid data format from AI model: {e}")
    except HTTPException as e:
        raise e
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate study plan: {str(e)}")


@router.post("/api/submit-answers")
async def submit_answers(request: AnswerSubmissionRequest, credentials: HTTPAuthorizationCredentials = Depends(security)):
    user_id = await get_user_id(credentials)
    if not user_id and request.user_id:
        user_id = request.user_id
    if not user_id:
        raise HTTPException(status_code=403, detail="Authentication required")

    with sqlite3.connect(DB_PATH) as conn:
        c = conn.cursor()
        for answer in request.answers:
            is_correct = answer['user_answer'] == answer['correct_answer']
            c.execute('''
                INSERT INTO user_question_results (user_id, subject, is_correct, timestamp)
                VALUES (?, ?, ?, ?)
            ''', (
                user_id,
                request.subject,
                is_correct,
                datetime.now().isoformat()
            ))
        conn.commit()

    return {"status": "Answers recorded"}

@router.post("/api/summarize")
async def summarize_text(request: ContentRequest, user_id: Optional[str] = Depends(get_user_id)):
    if not request.text.strip():
        raise HTTPException(status_code=400, detail="No text provided for summarization.")

    # Define tone-specific prompts
    tone_prompts = {
        "soft": "You are a friendly, encouraging educational assistant that provides very detailed but easy to understand, well-structured summaries. Use a warm, supportive tone and Markdown for formatting.",
        "medium": "You are a professional educational assistant that provides detailed, well-structured summaries. Use a balanced, informative tone and Markdown for formatting.",
        "strict": "You are a formal, precise educational assistant that provides concise, well-structured summaries. Use a professional, academic tone and Markdown for formatting."
    }
    
    system_prompt = tone_prompts.get(request.tone, tone_prompts["medium"])

    prompt_messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": f"""Summarize the following content in a systematic, well-structured format.
Use clear headings (using Markdown # or ##) and bullet points (using Markdown - or *) to organize the information logically.
Give all the concepts of content and don't skip any point or info from content. Keep it SHORT length, extra concise and detailed, but stay within limit, give as most info as possible is less words.dont be ultra straightforward, be like you are explaining it in simple language.Be ultra accurate. you should not give wrong answers.Also extract the key points and tell a conclusion in end.

Content to summarize:
{request.text.strip()}
"""}
    ]
    
    try:
        summary = call_llm(prompt_messages, max_tokens=1024)
        if user_id and request.text.strip():
            subject = detect_subject(request.text)
            with sqlite3.connect(DB_PATH) as conn:
                c = conn.cursor()
                c.execute('INSERT INTO user_studies (user_id, subject, mode, timestamp) VALUES (?, ?, ?, ?)',
                          (user_id, subject, 'summary', datetime.now().isoformat()))
                conn.commit()
        return {"summary": summary.strip()}
    except HTTPException as e:
        raise e
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate summary: {str(e)}")

@router.post("/api/generate-questions")
async def generate_questions(request: GenerateQuestionsRequest, user_id: Optional[str] = Depends(get_user_id)):
    if not request.text.strip():
        raise HTTPException(status_code=400, detail="No text provided for question generation.")

    prompt_content = f"""Generate {request.count} multiple-choice questions with a difficulty level of {request.difficulty} from the text below.
Each question must have 4 distinct options, and one correct answer.
THE QUESTIONS SHOULD BE EXACTLY FROM THE PROVIDED TEXT, NOT IRRELAVANT.
DONT ASK ULTRA BASIC QUESTIONS, THE LEVEL SHOULD BE HIGH.

Format your output strictly as a JSON list of objects. Do not include any other text or markdown outside the JSON.

Example JSON format:
[
  {{
    "text": "What is the capital of France?",
    "options": ["Paris", "London", "Rome", "Berlin"],
    "answer": "Paris"
  }},
  {{
    "text": "Which planet is known as the Red Planet?",
    "options": ["Mars", "Jupiter", "Venus", "Saturn"],
    "answer": "Mars"
  }}
]

Text to generate questions from:
{request.text.strip()}
"""
    prompt_messages = [
        {"role": "system", "content": "You are an educational assistant specialized in generating multiple-choice questions in strict JSON format."},
        {"role": "user", "content": prompt_content}
    ]

    try:
        raw_response = call_llm(prompt_messages, temperature=0.7, max_tokens=150 * request.count)
        if raw_response.startswith("```json") and raw_response.endswith("```"):
            raw_response = raw_response[len("```json"):-len("```")].strip()
        elif raw_response.startswith("```") and raw_response.endswith("```"):
            raw_response = raw_response[len("```"):-len("```")].strip()
        
        questions_data = json.loads(raw_response)
        if not isinstance(questions_data, list):
            raise HTTPException(status_code=500, detail="Model did not return a valid JSON list of questions.")
        
        final_questions = []
        for q in questions_data:
            if not all(k in q for k in ["text", "options", "answer"]):
                continue
            if not isinstance(q['options'], list) or len(q['options']) != 4:
                continue
            if q['answer'] not in q['options']:
                continue
            options = q['options']
            random.shuffle(options)
            final_questions.append({
                'text': q['text'],
                'options': options,
                'answer': q['answer']
            })
        
        if not final_questions:
            raise HTTPException(status_code=400, detail="No valid questions could be generated from the provided text.")
        
        return {"questions": final_questions}
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=500, detail=f"Failed to parse JSON from AI model: {e}")
    except HTTPException as e:
        raise e
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate questions: {str(e)}")


@router.post("/api/follow-up-question")
async def follow_up_question(request: FollowUpRequest, user_id: Optional[str] = Depends(get_user_id)):
    if not request.summary.strip() or not request.question.strip():
        raise HTTPException(status_code=400, detail="Summary or question missing for follow-up.")

    prompt_messages = [
        {"role": "system", "content": "You are a helpful educational assistant."},
        {"role": "user", "content": f"Based on the summary: {request.summary}\nAnswer the question: {request.question}"}
    ]
    
    try:
        answer = call_llm(prompt_messages, max_tokens=512)
        if user_id and request.summary.strip():
            subject = detect_subject(request.summary)
            with sqlite3.connect(DB_PATH) as conn:
                c = conn.cursor()
                c.execute('INSERT INTO user_studies (user_id, subject, mode, timestamp) VALUES (?, ?, ?, ?)',
                          (user_id, subject, 'follow-up', datetime.now().isoformat()))
                conn.commit()
        return {"answer": answer.strip()}
    except HTTPException as e:
        raise e
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to answer follow-up question: {str(e)}")

@router.post("/api/fetch-and-extract-url")
async def fetch_and_extract_url(request: URLRequest, user_id: Optional[str] = Depends(get_user_id)):
    if not request.url.strip():
        raise HTTPException(status_code=400, detail="No URL provided.")
    
    text = await extract_text_from_url(request.url)
    if user_id and text.strip():
        subject = detect_subject(text)
        with sqlite3.connect(DB_PATH) as conn:
            c = conn.cursor()
            c.execute('INSERT INTO user_studies (user_id, subject, mode, timestamp) VALUES (?, ?, ?, ?)',
                      (user_id, subject, 'url-extract', datetime.now().isoformat()))
            conn.commit()
    return {"text": text}

@router.post("/api/fetch-youtube-content")
async def fetch_youtube_content(request: YouTubeURLRequest, user_id: Optional[str] = Depends(get_user_id)):
    if not request.youtube_url.strip():
        raise HTTPException(status_code=400, detail="No YouTube URL provided.")
    
    transcript = extract_youtube_transcript(request.youtube_url)
    if user_id and transcript.strip():
        subject = detect_subject(transcript)
        with sqlite3.connect(DB_PATH) as conn:
            c = conn.cursor()
            c.execute('INSERT INTO user_studies (user_id, subject, mode, timestamp) VALUES (?, ?, ?, ?)',
                      (user_id, subject, 'youtube', datetime.now().isoformat()))
            conn.commit()
    return {"transcript": transcript}

@router.post("/api/generate-flashcards")
async def generate_flashcards_endpoint(request: FlashcardsRequest, user_id: Optional[str] = Depends(get_user_id)):
    if not request.text.strip():
        raise HTTPException(status_code=400, detail="No text provided for flashcards generation.")
    
    flashcards = generate_flashcards(request.text)
    if user_id and request.text.strip():
        subject = detect_subject(request.text)
        with sqlite3.connect(DB_PATH) as conn:
            c = conn.cursor()
            c.execute('INSERT INTO user_studies (user_id, subject, mode, timestamp) VALUES (?, ?, ?, ?)',
                      (user_id, subject, 'flashcards', datetime.now().isoformat()))
            conn.commit()
    return {"flashcards": flashcards}

@router.post("/api/generate-vocabulary")
async def generate_vocabulary_endpoint(request: VocabularyRequest, user_id: Optional[str] = Depends(get_user_id)):
    if not request.text.strip():
        raise HTTPException(status_code=400, detail="No text provided for vocabulary generation.")
    
    vocabulary = generate_vocabulary(request.text)
    if user_id and request.text.strip():
        subject = detect_subject(request.text)
        with sqlite3.connect(DB_PATH) as conn:
            c = conn.cursor()
            c.execute('INSERT INTO user_studies (user_id, subject, mode, timestamp) VALUES (?, ?, ?, ?)',
                      (user_id, subject, 'vocabulary', datetime.now().isoformat()))
            conn.commit()
    return {"vocabulary": vocabulary}

@router.post("/api/humanize-text")
async def humanize_text_endpoint(request: HumanizeRequest, user_id: Optional[str] = Depends(get_user_id)):
    if not request.text.strip():
        raise HTTPException(status_code=400, detail="No text provided for humanization.")
    
    humanized = humanize_text(request.text)
    if user_id and request.text.strip():
        subject = detect_subject(request.text)
        with sqlite3.connect(DB_PATH) as conn:
            c = conn.cursor()
            c.execute('INSERT INTO user_studies (user_id, subject, mode, timestamp) VALUES (?, ?, ?, ?)',
                      (user_id, subject, 'humanize', datetime.now().isoformat()))
            conn.commit()
    return {"humanized_text": humanized}

@router.post("/api/correct-grammar")
async def correct_grammar_endpoint(request: GrammarCorrectionRequest, user_id: Optional[str] = Depends(get_user_id)):
    if not request.text.strip():
        raise HTTPException(status_code=400, detail="No text provided for grammar correction.")
    
    corrected = correct_grammar(request.text)
    if user_id and request.text.strip():
        subject = detect_subject(request.text)
        with sqlite3.connect(DB_PATH) as conn:
            c = conn.cursor()
            c.execute('INSERT INTO user_studies (user_id, subject, mode, timestamp) VALUES (?, ?, ?, ?)',
                      (user_id, subject, 'grammar-correction', datetime.now().isoformat()))
            conn.commit()
    return {"corrected_text": corrected}

@router.post("/api/generate-mindmap")
async def generate_mindmap_endpoint(request: MindMapRequest, user_id: Optional[str] = Depends(get_user_id)):
    if not request.text.strip():
        raise HTTPException(status_code=400, detail="No text provided for mindmap generation.")
    
    mindmap_code = generate_mindmap(request.text)
    if user_id and request.text.strip():
        subject = detect_subject(request.text)
        with sqlite3.connect(DB_PATH) as conn:
            c = conn.cursor()
            c.execute('INSERT INTO user_studies (user_id, subject, mode, timestamp) VALUES (?, ?, ?, ?)',
                      (user_id, subject, 'mindmap', datetime.now().isoformat()))
            conn.commit()
    return {"mindmap_code": mindmap_code}

@router.post("/api/generate-diagram")
async def generate_diagram_endpoint(request: DiagramRequest, user_id: Optional[str] = Depends(get_user_id)):
    if not request.text.strip():
        raise HTTPException(status_code=400, detail="No text provided for diagram generation.")
    
    diagram_code = generate_diagram(request.text, request.diagram_type)
    if user_id and request.text.strip():
        subject = detect_subject(request.text)
        with sqlite3.connect(DB_PATH) as conn:
            c = conn.cursor()
            c.execute('INSERT INTO user_studies (user_id, subject, mode, timestamp) VALUES (?, ?, ?, ?)',
                      (user_id, subject, 'diagram', datetime.now().isoformat()))
            conn.commit()
    return {"diagram_code": diagram_code}

@router.post("/api/generate-handwritten")
async def generate_handwritten_endpoint(request: HandwrittenRequest, user_id: Optional[str] = Depends(get_user_id)):
    if not request.text.strip():
        raise HTTPException(status_code=400, detail="No text provided for handwritten notes generation.")
    
    handwritten_text = generate_handwritten(request.text, request.style)
    if user_id and request.text.strip():
        subject = detect_subject(request.text)
        with sqlite3.connect(DB_PATH) as conn:
            c = conn.cursor()
            c.execute('INSERT INTO user_studies (user_id, subject, mode, timestamp) VALUES (?, ?, ?, ?)',
                      (user_id, subject, 'handwritten', datetime.now().isoformat()))
            conn.commit()
    return {"handwritten_text": handwritten_text}

class SolveQuestionsRequest(BaseModel):
    text: str
    tone: str = "medium"  # Add tone parameter

# Update the solve_questions function to accept tone parameter
def solve_questions(text: str, tone: str = "medium") -> List[dict]:
    # Define tone-specific prompts
    tone_prompts = {
        "soft": "You are a friendly, encouraging educational assistant specialized in solving questions. Use a warm, supportive tone.",
        "medium": "You are a professional educational assistant specialized in solving questions. Use a balanced, informative tone.",
        "strict": "You are a formal, precise educational assistant specialized in solving questions. Use a professional, academic tone."
    }
    
    system_prompt = tone_prompts.get(tone, tone_prompts["medium"])

    prompt_content = f"""Solve the questions in the following content.
For each question, provide a clear answer and explanation.
Format the output strictly as a JSON object with a single key "answers", whose value is a list of answer objects.
Each answer object should have 'question' and 'answer'.
Do not include any other text or markdown outside the JSON.

Content with questions:
{text.strip()}
"""
    prompt_messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": prompt_content}
    ]

    try:
        raw_response = call_llm(prompt_messages, temperature=0.4, max_tokens=1500)
        if raw_response.startswith("```json") and raw_response.endswith("```"):
            raw_response = raw_response[len("```json"):-len("```")].strip()
        elif raw_response.startswith("```") and raw_response.endswith("```"):
            raw_response = raw_response[len("```"):-len("```")].strip()
        
        answers_data = json.loads(raw_response)
        if not isinstance(answers_data, dict) or "answers" not in answers_data or not isinstance(answers_data["answers"], list):
            raise ValueError("Model did not return a valid JSON object with an 'answers' list.")
        
        final_answers = []
        for a in answers_data["answers"]:
            if not all(k in a for k in ["question", "answer"]):
                continue
            final_answers.append({
                "question": a["question"].strip(),
                "answer": a["answer"].strip()
            })
        
        if not final_answers:
            raise HTTPException(status_code=400, detail="No valid answers could be generated from the provided content.")
        
        return final_answers
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=500, detail=f"Failed to parse JSON from AI model: {e}")
    except ValueError as e:
        raise HTTPException(status_code=500, detail=f"Invalid data format from AI model: {e}")
    except HTTPException as e:
        raise e
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to solve questions: {str(e)}")

class QuestionPaperRequest(BaseModel):
    text: str
    subject: str = "General Studies"
    difficulty: str = "medium"
    sections: Optional[List[str]] = ["MCQ", "Short Answer", "Long Answer"]
    format: str = "pdf"  # "pdf" or "docx"

from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from docx import Document

@router.post("/api/generate-question-paper")
async def generate_question_paper(request: QuestionPaperRequest, user_id: Optional[str] = Depends(get_user_id)):
    if not request.text.strip():
        raise HTTPException(status_code=400, detail="No content provided")

    prompt = f"""
    Generate a structured question paper for subject: {request.text}.
    Difficulty: {request.difficulty}.
    Sections: {', '.join(request.sections)}.

    Format it like a real exam:
    - Title and Instructions
    - Section A: Multiple Choice
    - Section B: Short Answer
    - Section C: Long Answer
    Number the questions clearly.
    """

    paper_text = call_llm([{"role": "user", "content": prompt}], temperature=0.6, max_tokens=2000)

    if request.format == "pdf":
        filename = f"/tmp/question_paper_{uuid.uuid4().hex}.pdf"
        c = canvas.Canvas(filename, pagesize=A4)
        width, height = A4
        y = height - 50

        for line in paper_text.split("\n"):
            if not line.strip():
                y -= 15
                continue
            c.drawString(50, y, line.strip())
            y -= 20
            if y < 50:
                c.showPage()
                y = height - 50
        c.save()
        return FileResponse(filename, filename="QuestionPaper.pdf", media_type="application/pdf")

    elif request.format == "docx":
        filename = f"/tmp/question_paper_{uuid.uuid4().hex}.docx"
        doc = Document()
        for line in paper_text.split("\n"):
            doc.add_paragraph(line.strip())
        doc.save(filename)
        return FileResponse(filename, filename="QuestionPaper.docx", media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

    else:
        raise HTTPException(status_code=400, detail="Invalid format. Use 'pdf' or 'docx'.")

# Update the solve_questions_endpoint function
@router.post("/api/solve-questions")
async def solve_questions_endpoint(request: SolveQuestionsRequest, user_id: Optional[str] = Depends(get_user_id)):
    if not request.text.strip():
        raise HTTPException(status_code=400, detail="No questions provided to solve.")
    
    answers = solve_questions(request.text, request.tone)
    if user_id and request.text.strip():
        subject = detect_subject(request.text)
        with sqlite3.connect(DB_PATH) as conn:
            c = conn.cursor()
            c.execute('INSERT INTO user_studies (user_id, subject, mode, timestamp) VALUES (?, ?, ?, ?)',
                      (user_id, subject, 'solve-questions', datetime.now().isoformat()))
            conn.commit()
    return {"answers": answers}

from fastapi import UploadFile, File
import whisper
import tempfile

@router.post("/api/upload-audio")
async def upload_audio(file: UploadFile = File(...), user_id: Optional[str] = Depends(get_user_id)):
    # Accept more audio formats
    os.environ["HF_HOME"] = "/tmp/hf_cache"
    os.makedirs("/tmp/hf_cache", exist_ok=True)
    allowed_content_types = [
        'audio/mpeg', 'audio/mp3', 'audio/wav', 'audio/wave', 
        'audio/x-wav', 'audio/webm', 'audio/ogg', 'audio/m4a',
        'audio/x-m4a', 'audio/aac'
    ]
    
    if file.content_type not in allowed_content_types:
        raise HTTPException(status_code=400, detail=f"Unsupported audio format. Supported formats: MP3, WAV, WebM, OGG, M4A, AAC")
    
    try:
        # Save temp file
        with tempfile.NamedTemporaryFile(delete=False, suffix=".audio") as tmp:
            content = await file.read()
            tmp.write(content)
            tmp_path = tmp.name

        # Transcribe with Whisper
        model = whisper.load_model("small")
        result = model.transcribe(tmp_path)
        
        # Clean up temp file
        os.unlink(tmp_path)

        text = result["text"]

        if not text.strip():
            raise HTTPException(status_code=400, detail="No speech detected in audio")

        if len(text) > MAX_LLM_INPUT_CHARS:
            text = text[:MAX_LLM_INPUT_CHARS]

        return {"extracted_text": text}

    except Exception as e:
        # Clean up temp file if it exists
        if 'tmp_path' in locals() and os.path.exists(tmp_path):
            os.unlink(tmp_path)
        raise HTTPException(status_code=500, detail=f"Failed to transcribe audio: {str(e)}")
class DoubtsRequest(BaseModel):
    text: str
def generate_doubts(text: str) -> List[dict]:
    prompt_content = f"""Generate several potential doubts a student might have while studying the following content.
Each doubt should include a 'question' (the doubt itself), 'context' (why this doubt might arise), and 'hint' (a helpful tip to resolve the doubt).
Format the output strictly as a JSON object with a single key "doubts", whose value is a list of doubt objects.
Each doubt object should have 'question', 'context', and 'hint'.
Do not include any other text or markdown outside the JSON.

Content:
{text.strip()}
"""
    prompt_messages = [
        {"role": "system", "content": "You are an educational assistant specialized in identifying potential student doubts in strict JSON format."},
        {"role": "user", "content": prompt_content}
    ]

    try:
        raw_response = call_llm(prompt_messages, temperature=0.6, max_tokens=1000)
        if raw_response.startswith("```json") and raw_response.endswith("```"):
            raw_response = raw_response[len("```json"):-len("```")].strip()
        elif raw_response.startswith("```") and raw_response.endswith("```"):
            raw_response = raw_response[len("```"):-len("```")].strip()
        
        doubts_data = json.loads(raw_response)
        if not isinstance(doubts_data, dict) or "doubts" not in doubts_data or not isinstance(doubts_data["doubts"], list):
            raise ValueError("Model did not return a valid JSON object with a 'doubts' list.")
        
        final_doubts = []
        for doubt in doubts_data["doubts"]:
            if not all(k in doubt for k in ["question", "context", "hint"]):
                continue
            final_doubts.append({
                "question": doubt["question"].strip(),
                "context": doubt["context"].strip(),
                "hint": doubt["hint"].strip()
            })
        
        if not final_doubts:
            raise HTTPException(status_code=400, detail="No valid doubts could be generated from the provided content.")
        
        return final_doubts
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=500, detail=f"Failed to parse JSON from AI model: {e}")
    except ValueError as e:
        raise HTTPException(status_code=500, detail=f"Invalid data format from AI model: {e}")
    except HTTPException as e:
        raise e
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate doubts: {str(e)}")

# New Doubts Endpoint
@router.post("/api/generate-doubts")
async def generate_doubts_endpoint(request: DoubtsRequest, user_id: Optional[str] = Depends(get_user_id)):
    if not request.text.strip():
        raise HTTPException(status_code=400, detail="No text provided for doubts generation.")
    
    doubts = generate_doubts(request.text)
    if user_id and request.text.strip():
        subject = detect_subject(request.text)
        with sqlite3.connect(DB_PATH) as conn:
            c = conn.cursor()
            c.execute('INSERT INTO user_studies (user_id, subject, mode, timestamp) VALUES (?, ?, ?, ?)',
                      (user_id, subject, 'doubts', datetime.now().isoformat()))
            conn.commit()
    return {"doubts": doubts}
@router.post("/api/generate-ppt")
async def generate_ppt_endpoint(request: PPTRequest, user_id: Optional[str] = Depends(get_user_id)):
    if not request.text.strip():
        raise HTTPException(status_code=400, detail="No text provided for PPT generation.")
    
    ppt_path = generate_ppt(request.text, request.slide_count, request.font_family, request.font_color, request.slide_bg_color, request.font_size)
    if user_id and request.text.strip():
        subject = detect_subject(request.text)
        with sqlite3.connect(DB_PATH) as conn:
            c = conn.cursor()
            c.execute('INSERT INTO user_studies (user_id, subject, mode, timestamp) VALUES (?, ?, ?, ?)',
                      (user_id, subject, 'ppt', datetime.now().isoformat()))
            conn.commit()
    return FileResponse(
        path=ppt_path,
        filename=f"presentation-{int(time.time())}.pptx",
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

# Add to imports
from fastapi.responses import JSONResponse

# Add to Pydantic Models
class AIDebateRequest(BaseModel):
    topic: str
    user_position: str
    debate_rounds: int = 3

# Add to Feature Costs


# Add to Helper Functions
def generate_debate(topic: str, user_position: str, rounds: int) -> List[dict]:
    prompt_content = f"""Simulate a debate on the topic: '{topic}'.
    - There are three participants: User, AI_Persona_1, and AI_Persona_2.
    - User position: {user_position}.
    - AI_Persona_1 takes a supportive stance, AI_Persona_2 takes an opposing stance.
    - Generate {rounds} rounds of debate.
    - Each round includes: User's argument, AI_Persona_1's supportive response, AI_Persona_2's counterargument.
    - Format output as JSON with a single key "debate_rounds", containing a list of objects with keys: 'round', 'user_argument', 'ai_persona_1_response', 'ai_persona_2_response'.
    - Ensure responses are concise, relevant, and maintain distinct POVs.
    - Do not include any text or markdown outside the JSON.
    """
    prompt_messages = [
        {"role": "system", "content": "You are an AI simulating a structured debate with distinct perspectives in strict JSON format."},
        {"role": "user", "content": prompt_content}
    ]

    try:
        raw_response = call_llm(prompt_messages, temperature=0.7, max_tokens=2000)
        if raw_response.startswith("```json") and raw_response.endswith("```"):
            raw_response = raw_response[len("```json"):-len("```")].strip()
        elif raw_response.startswith("```") and raw_response.endswith("```"):
            raw_response = raw_response[len("```"):-len("```")].strip()
        
        debate_data = json.loads(raw_response)
        if not isinstance(debate_data, dict) or "debate_rounds" not in debate_data or not isinstance(debate_data["debate_rounds"], list):
            raise ValueError("Model did not return a valid JSON object with a 'debate_rounds' list.")
        
        final_debate = []
        for round_data in debate_data["debate_rounds"][:rounds]:
            if not all(k in round_data for k in ["round", "user_argument", "ai_persona_1_response", "ai_persona_2_response"]):
                continue
            final_debate.append({
                "round": round_data["round"],
                "user_argument": round_data["user_argument"].strip(),
                "ai_persona_1_response": round_data["ai_persona_1_response"].strip(),
                "ai_persona_2_response": round_data["ai_persona_2_response"].strip()
            })
        
        if not final_debate:
            raise HTTPException(status_code=400, detail="No valid debate rounds could be generated.")
        
        return final_debate
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=500, detail=f"Failed to parse JSON from AI model: {e}")
    except ValueError as e:
        raise HTTPException(status_code=500, detail=f"Invalid data format from AI model: {e}")
    except HTTPException as e:
        raise e
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate debate: {str(e)}")


@router.post("/api/generate-ai-debate")
async def generate_ai_debate_endpoint(request: AIDebateRequest, user_id: Optional[str] = Depends(get_user_id)):
    if not request.topic.strip():
        raise HTTPException(status_code=400, detail="No topic provided for debate generation.")
    
    debate_rounds = generate_debate(request.topic, request.user_position, request.debate_rounds)
    
    if user_id and request.topic.strip():
        subject = detect_subject(request.topic)
        with sqlite3.connect(DB_PATH) as conn:
            c = conn.cursor()
            c.execute('INSERT INTO user_studies (user_id, subject, mode, timestamp) VALUES (?, ?, ?, ?)',
                      (user_id, subject, 'ai_debate', datetime.now().isoformat()))
            conn.commit()
    
    return JSONResponse(content={"debate_rounds": debate_rounds})
@router.post("/api/generate-study-plan")
async def generate_study_plan_endpoint(request: StudyPlanRequest, user_id: Optional[str] = Depends(get_user_id)):
    if not request.text.strip():
        raise HTTPException(status_code=400, detail="No text provided for study plan generation.")
    
    study_plan = generate_study_plan(request)
    if user_id and request.text.strip():
        subject = detect_subject(request.text)
        with sqlite3.connect(DB_PATH) as conn:
            c = conn.cursor()
            c.execute('INSERT INTO user_studies (user_id, subject, mode, timestamp) VALUES (?, ?, ?, ?)',
                      (user_id, subject, 'study-plan', datetime.now().isoformat()))
            conn.commit()
    return study_plan

@router.post("/api/process-file")
async def process_file(file: UploadFile = File(...), user_id: Optional[str] = Depends(get_user_id)):
    content_type = file.content_type
    if content_type == 'application/pdf':
        text = await extract_text_from_pdf(file)
    elif content_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
        text = await extract_text_from_docx(file)
    elif content_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation', 'application/vnd.ms-powerpoint']:
        text = await extract_text_from_ppt(file)
    elif content_type.startswith('image/'):
        text = await extract_text_from_image(file)
    else:
        raise HTTPException(status_code=400, detail="Unsupported file type.")
    
    if not text.strip():
        raise HTTPException(status_code=400, detail="No text could be extracted from the file.")
    
    if user_id and text.strip():
        subject = detect_subject(text)
        with sqlite3.connect(DB_PATH) as conn:
            c = conn.cursor()
            c.execute('INSERT INTO user_studies (user_id, subject, mode, timestamp) VALUES (?, ?, ?, ?)',
                      (user_id, subject, 'file-process', datetime.now().isoformat()))
            conn.commit()
    return {"text": text}

# @router.post("/api/generate-referral")
# async def generate_referral(request: ReferralRequest):
#     # ... (Existing referral logic remains unchanged)
#     pass

# @router.post("/api/validate-referral")
# async def validate_referral(request: ReferralValidationRequest):
#     # ... (Existing referral validation logic remains unchanged)
#     pass

# @router.post("/api/referral-status")
# async def referral_status(request: ReferralStatusRequest):
#     # ... (Existing referral status logic remains unchanged)
#     pass

@router.post("/api/guest-usage")
async def update_guest_usage(request: GuestUsageRequest):
    # ... (Existing guest usage logic remains unchanged)
    pass

def detect_subject(text: str) -> str:
    prompt = f"""Classify the following text into one main subject from this list: Math, Physics, Chemistry, Biology, History, Geography, Literature, Computer Science, Economics, Art, Other.
Output ONLY the subject name, nothing else.

Text:
{text[:500]}"""  # Limit to first 500 chars to save tokens
    
    prompt_messages = [{"role": "user", "content": prompt}]
    try:
        subject = call_llm(prompt_messages, temperature=0.2, max_tokens=10).strip()
        return subject if subject in ['Math', 'Physics', 'Chemistry', 'Biology', 'History', 'Geography', 'Literature', 'Computer Science', 'Economics', 'Art', 'Other'] else 'Other'
    except Exception:
        return 'Other'  # Fallback

@router.get("/api/dashboard")
async def get_dashboard(user_id: Optional[str] = Depends(get_user_id)):
    if not user_id:
        raise HTTPException(status_code=403, detail="Authentication required for dashboard access")
    
    with sqlite3.connect(DB_PATH) as conn:
        c = conn.cursor()
        
        # Get studied subjects
        c.execute('''
            SELECT subject, COUNT(*) as value
            FROM user_studies
            WHERE user_id = ?
            GROUP BY subject
            ORDER BY value DESC
        ''', (user_id,))
        subjects = [{"name": row[0], "value": row[1]} for row in c.fetchall()]
        
        # Get subject performance (correct vs incorrect answers)
        c.execute('''
            SELECT subject, 
                   SUM(CASE WHEN is_correct = 1 THEN 1 ELSE 0 END) as correct,
                   SUM(CASE WHEN is_correct = 0 THEN 1 ELSE 0 END) as incorrect
            FROM user_question_results
            WHERE user_id = ?
            GROUP BY subject
        ''', (user_id,))
        
        subject_performance = []
        for row in c.fetchall():
            subject, correct, incorrect = row
            total = correct + incorrect
            accuracy = round((correct / total) * 100, 2) if total > 0 else 0
            subject_performance.append({
                "subject": subject,
                "correct": correct,
                "incorrect": incorrect,
                "accuracy": accuracy
            })
        
        # Sort by accuracy to determine strong/weak subjects
        subject_performance.sort(key=lambda x: x["accuracy"], reverse=True)
        
        strong_subjects = subject_performance[:3]  # Top 3
        weak_subjects = subject_performance[-3:]   # Bottom 3
    
    return {
        "subjects": subjects,
        "strong_subjects": strong_subjects,
        "weak_subjects": weak_subjects
    }

# --- Blog Endpoints ---
@router.get("/api/blog/post/{post_id}")
async def get_blog_post(post_id: str):
    try:
        logger.info(f"Attempting to fetch post: {post_id}")
        
        # Test Supabase connection first
        try:
            # Simple query to test connection
            test_response = supabase.table("blog_posts").select("count").execute()
            logger.info(f"Supabase connection test: {test_response}")
        except Exception as conn_error:
            logger.error(f"Supabase connection failed: {conn_error}")
            raise HTTPException(status_code=500, detail="Database connection failed")
        
        # Fetch the specific post
        response = supabase.table("blog_posts").select("*").eq("id", post_id).execute()
        logger.info(f"Raw response: {response}")
        
        if not response.data:
            logger.warning(f"Post not found: {post_id}")
            raise HTTPException(status_code=404, detail="Post not found")
        
        if len(response.data) > 1:
            logger.warning(f"Multiple posts found with same ID: {post_id}")
        
        post = response.data[0]
        logger.info(f"Post data: {post}")
        
        # Process tags
        tags_value = post.get('tags')
        if isinstance(tags_value, str) and tags_value.strip():
            try:
                post['tags'] = json.loads(tags_value)
            except json.JSONDecodeError:
                logger.warning(f"Invalid JSON in tags for post {post_id}: {tags_value}")
                post['tags'] = []
        else:
            post['tags'] = []
        
        return post
        
    except HTTPException as he:
        raise he
    except Exception as e:
        logger.error(f"Unexpected error fetching post {post_id}: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail="Internal server error")

@router.post("/api/blog/posts")
async def create_blog_post(
    post: BlogPostCreate,
    credentials: HTTPAuthorizationCredentials = Depends(security)
):
    user_id = await get_user_id(credentials)
    if not user_id:
        raise HTTPException(status_code=403, detail="Authentication required")
    
    # Generate unique ID
    post_id = str(uuid.uuid4())
    
    # Supabase mein insert karo
    response = supabase.table("blog_posts").insert({
        "id": post_id,
        "title": post.title,
        "content": post.content,
        "author": post.author,
        "tags": json.dumps(post.tags),
        "created_at": datetime.now().isoformat(),
        "updated_at": datetime.now().isoformat()
    }).execute()
    
    return {"id": post_id, "status": "created"}

@router.put("/api/blog/post/{post_id}")
async def update_blog_post(
    post_id: str,
    post: BlogPostCreate,
    credentials: HTTPAuthorizationCredentials = Depends(security)
):
    user_id = await get_user_id(credentials)
    if not user_id:
        raise HTTPException(status_code=403, detail="Authentication required")
    
    # Supabase mein update karo
    response = supabase.table("blog_posts").update({
        "title": post.title,
        "content": post.content,
        "author": post.author,
        "tags": json.dumps(post.tags),
        "updated_at": datetime.now().isoformat()
    }).eq("id", post_id).execute()
    
    if not response.data:
        raise HTTPException(status_code=404, detail="Post not found")
    
    return {"status": "updated"}

@router.delete("/api/blog/post/{post_id}")
async def delete_blog_post(
    post_id: str,
    credentials: HTTPAuthorizationCredentials = Depends(security)
):
    user_id = await get_user_id(credentials)
    if not user_id:
        raise HTTPException(status_code=403, detail="Authentication required")
    
    # Supabase se delete karo
    response = supabase.table("blog_posts").delete().eq("id", post_id).execute()
    
    if not response.data:
        raise HTTPException(status_code=404, detail="Post not found")
    
    return {"status": "deleted"}

# ---------- Database Initialization (SQLite Blog Table Removal) ----------


# Add to imports


MAX_LLM_INPUT_CHARS = 28000

DB_PATH = "/tmp/nexnotes.db"

# --- Pydantic Models ---
class Mention(BaseModel):
    doc_id: str
    text: str
    span: tuple[int, int]

class Concept(BaseModel):
    name: str
    type: str
    score: float
    mentions: List[Mention]

class ConceptsResponse(BaseModel):
    concepts: List[Concept]

class KnowledgeGraphNode(BaseModel):
    id: str
    label: str
    type: str
    sources: Dict[str, List[str]]
    confidence: float

class KnowledgeGraphEdge(BaseModel):
    source: str
    target: str
    relationship: str
    sources: Dict[str, List[str]]

class KnowledgeGraphResponse(BaseModel):
    nodes: List[KnowledgeGraphNode]
    edges: List[KnowledgeGraphEdge]
    conflicts: List[Dict]
    gaps: List[Dict]

class KnowledgeGraphRequest(BaseModel):
    documents: List[str]
    course_name: str = "General Knowledge"

# --- Helper Functions ---


# --- Concept Extraction ---
def _concepts_system_prompt() -> Dict[str, str]:
    return {
        "role": "system",
        "content": (
            "You are a strict JSON generator extracting structured knowledge from documents. "
            "Return ONLY valid JSON matching this schema: "
            '{"concepts": [{"name": "string", "type": "string", "score": number, "mentions": [{"doc_id": "string", "text": "string", "span": [integer, integer]}]}]}'
            "Do not include prose, markdown, or code fences. Ensure the output is valid JSON."
        )
    }

def _concepts_user_prompt(doc_summaries: List[Dict[str, str]]) -> Dict[str, str]:
    payload = {
        "schema": {
            "type": "object",
            "properties": {
                "concepts": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "required": ["name", "type", "score", "mentions"],
                        "properties": {
                            "name": {"type": "string", "description": "Name of the concept"},
                            "type": {"type": "string", "enum": ["concept", "definition", "example"]},
                            "score": {"type": "number", "description": "Confidence score between 0 and 1"},
                            "mentions": {
                                "type": "array",
                                "items": {
                                    "type": "object",
                                    "required": ["doc_id", "text", "span"],
                                    "properties": {
                                        "doc_id": {"type": "string", "description": "Document identifier"},
                                        "text": {"type": "string", "description": "Text of the mention"},
                                        "span": {
                                            "type": "array",
                                            "items": {"type": "integer"},
                                            "minItems": 2,
                                            "maxItems": 2,
                                            "description": "Start and end character indices"
                                        }
                                    }
                                }
                            }
                        }
                    }
                }
            },
            "required": ["concepts"]
        },
        "documents": [
            {"doc_id": doc["doc_id"], "text": doc["text"][:5000]}
            for doc in doc_summaries
        ]
    }
    return {"role": "user", "content": json.dumps(payload, ensure_ascii=False)}

def _parse_llm_json(raw: str) -> Dict:
    if not raw.strip():
        logger.error("Empty response from LLM")
        raise ValueError("Empty response from LLM")
    try:
        return json.loads(raw)
    except json.JSONDecodeError as e:
        logger.error(f"Failed to parse JSON: {str(e)}")
        logger.error(f"Raw response: {raw[:1000]}...")
        raise

def extract_concepts_from_documents(
    documents: List[Dict[str, Any]],
    call_llm_fn=call_llm,
    temperature: float = 0.0,
    max_tokens: int = 3500,
    retries: int = 3
) -> ConceptsResponse:
    doc_summaries: List[Dict[str, str]] = []
    for i, d in enumerate(documents):
        if isinstance(d, str):
            doc_summaries.append({"doc_id": f"doc_{i+1}", "text": d.strip()})
        elif isinstance(d, dict):
            doc_id = d.get("doc_id") or d.get("id") or d.get("name") or f"doc_{i+1}"
            text = d.get("text") or d.get("content") or ""
            doc_summaries.append({"doc_id": str(doc_id), "text": str(text).strip()})
        else:
            logger.error(f"Unsupported document type: {type(d).__name__}")
            raise HTTPException(status_code=400, detail=f"Unsupported document type: {type(d).__name__}")

    total_chars = sum(len(doc["text"]) for doc in doc_summaries)
    if total_chars > MAX_LLM_INPUT_CHARS:
        logger.warning(f"Total input size {total_chars} exceeds {MAX_LLM_INPUT_CHARS}, truncating")
        for doc in doc_summaries:
            doc["text"] = doc["text"][:4000]

    messages = [
        _concepts_system_prompt(),
        _concepts_user_prompt(doc_summaries)
    ]

    fallback_messages = [
        {
            "role": "system",
            "content": (
                "You are a strict JSON generator. Return ONLY valid JSON: "
                '{"concepts": [{"name": "string", "type": "string", "score": 0.5, "mentions": [{"doc_id": "string", "text": "string", "span": [0, 0]}]}]}'
            )
        },
        {
            "role": "user",
            "content": json.dumps({
                "documents": [
                    {"doc_id": doc["doc_id"], "text": doc["text"][:2000]}
                    for doc in doc_summaries
                ]
            }, ensure_ascii=False)
        }
    ]

    last_err: Optional[Exception] = None
    for attempt in range(retries):
        try:
            logger.info(f"Attempt {attempt + 1}: Sending prompt with {len(json.dumps(messages))} chars")
            raw = call_llm_fn(messages, temperature=temperature, max_tokens=max_tokens)
            logger.info(f"LLM Response: {raw[:500]}...")
            
            parsed = _parse_llm_json(raw)
            resp = ConceptsResponse(**parsed)

            clean_concepts: List[Concept] = []
            for c in resp.concepts[:25]:
                score = min(max(c.score, 0.0), 1.0)
                fixed_mentions: List[Mention] = []
                for m in c.mentions:
                    s0, s1 = int(m.span[0]), int(m.span[1])
                    fixed_mentions.append(Mention(doc_id=m.doc_id, text=m.text, span=(s0, s1)))
                clean_concepts.append(Concept(name=c.name.strip(), type=c.type.strip(), score=score, mentions=fixed_mentions))

            return ConceptsResponse(concepts=clean_concepts)

        except (json.JSONDecodeError, ValidationError, ValueError) as e:
            last_err = e
            logger.warning(f"Concepts parse attempt {attempt + 1} failed: {str(e)}")
            if attempt < retries - 1:
                logger.info("Switching to fallback prompt")
                messages = fallback_messages
            continue
        except Exception as e:
            last_err = e
            logger.exception(f"LLM call failed on attempt {attempt + 1}: {str(e)}")
            continue

    logger.warning("All attempts failed, returning minimal fallback response")
    return ConceptsResponse(concepts=[
        Concept(
            name="unknown",
            type="concept",
            score=0.5,
            mentions=[Mention(doc_id=doc_summaries[0]["doc_id"], text="Unknown concept", span=(0, 0))]
        )
    ] if doc_summaries else [])

# --- Knowledge Graph Endpoint ---
@router.post("/api/generate-knowledge-graph")
async def generate_knowledge_graph(
    request: KnowledgeGraphRequest,
    user_id: Optional[str] = Depends(get_user_id)
):
    if not request.documents or not any(doc.strip() for doc in request.documents):
        logger.error("No valid documents provided to generate_knowledge_graph")
        raise HTTPException(status_code=400, detail="No valid documents provided")
    
    try:
        doc_dicts = [{"doc_id": f"doc_{i+1}", "text": doc} for i, doc in enumerate(request.documents)]
        
        concepts_response = extract_concepts_from_documents(
            documents=doc_dicts,
            call_llm_fn=call_llm,
            temperature=0.0,
            max_tokens=3500,
            retries=3
        )
        
        concepts_dict = {
            "concepts": {
                concept.name: {
                    "definitions": {
                        mention.doc_id: mention.text
                        for mention in concept.mentions
                        if concept.type == "definition"
                    },
                    "examples": {
                        mention.doc_id: [mention.text]
                        for mention in concept.mentions
                        if concept.type == "example"
                    },
                    "relationships": [
                        {
                            "type": "related_to",
                            "target": other_concept.name,
                            "sources": [m.doc_id for m in concept.mentions]
                        }
                        for other_concept in concepts_response.concepts
                        if other_concept.name != concept.name
                        and any(m.doc_id in [m2.doc_id for m2 in other_concept.mentions] for m in concept.mentions)
                    ]
                }
                for concept in concepts_response.concepts
            },
            "formulas": {}
        }
        
        graph = build_knowledge_graph(concepts_dict, request.documents, request.course_name)
        
        analysis = analyze_knowledge_graph(graph)
        
        mermaid_code = convert_to_mermaid(graph)
        
        if user_id:
            with sqlite3.connect(DB_PATH) as conn:
                c = conn.cursor()
                c.execute('INSERT INTO user_studies (user_id, subject, mode, timestamp) VALUES (?, ?, ?, ?)',
                         (user_id, request.course_name, 'knowledge-graph', datetime.now().isoformat()))
                conn.commit()
        
        return {
            "graph": graph,
            "mermaid_code": mermaid_code,
            "conflicts": analysis["conflicts"],
            "gaps": analysis["gaps"]
        }
    except Exception as e:
        logger.error(f"Error in generate_knowledge_graph: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to generate knowledge graph: {str(e)}")

def build_knowledge_graph(concepts: Dict, documents: List[str], course_name: str) -> KnowledgeGraphResponse:
    nodes = []
    edges = []
    
    for concept, data in concepts.get("concepts", {}).items():
        node_id = f"concept_{len(nodes)}"
        nodes.append(KnowledgeGraphNode(
            id=node_id,
            label=concept,
            type="concept",
            sources=data.get("definitions", {}),
            confidence=calculate_confidence(data)
        ))
        
        for doc, definition in data.get("definitions", {}).items():
            def_id = f"def_{len(nodes)}"
            nodes.append(KnowledgeGraphNode(
                id=def_id,
                label=definition,
                type="definition",
                sources={doc: [f"Definition of {concept}"]},
                confidence=1.0
            ))
            edges.append(KnowledgeGraphEdge(
                source=node_id,
                target=def_id,
                relationship="defined_as",
                sources={doc: [f"Definition of {concept}"]}
            ))
        
        for rel in data.get("relationships", []):
            target_concept = rel["target"]
            if target_concept in concepts["concepts"]:
                edges.append(KnowledgeGraphEdge(
                    source=node_id,
                    target=f"concept_{list(concepts['concepts'].keys()).index(target_concept)}",
                    relationship=rel["type"],
                    sources={doc: [f"Relationship between {concept} and {target_concept}"] for doc in rel["sources"]}
                ))
    
    for formula, data in concepts.get("formulas", {}).items():
        formula_id = f"formula_{len(nodes)}"
        nodes.append(KnowledgeGraphNode(
            id=formula_id,
            label=data["expression"],
            type="formula",
            sources={doc: [f"Formula {formula}"] for doc in data["sources"]},
            confidence=1.0
        ))
    
    return KnowledgeGraphResponse(nodes=nodes, edges=edges, conflicts=[], gaps=[])

def analyze_knowledge_graph(graph: KnowledgeGraphResponse) -> Dict:
    conflicts = []
    gaps = []
    
    definition_nodes = [n for n in graph.nodes if n.type == "definition"]
    grouped_definitions = defaultdict(list)
    for node in definition_nodes:
        grouped_definitions[node.label.split(':')[0]].append(node)
    
    for concept, def_nodes in grouped_definitions.items():
        if len(def_nodes) > 1:
            unique_defs = len(set(n.label for n in def_nodes))
            if unique_defs > 1:
                conflicts.append({
                    "type": "definition_conflict",
                    "concept": concept,
                    "definitions": [n.label for n in def_nodes],
                    "sources": [list(n.sources.keys())[0] for n in def_nodes]
                })
    
    concept_nodes = [n for n in graph.nodes if n.type == "concept"]
    for node in concept_nodes:
        has_definition = any(
            e.relationship == "defined_as" and e.source == node.id 
            for e in graph.edges
        )
        if not has_definition:
            gaps.append({
                "type": "missing_definition",
                "concept": node.label,
                "sources": list(node.sources.keys())
            })
    
    return {"conflicts": conflicts, "gaps": gaps}

def convert_to_mermaid(graph: KnowledgeGraphResponse) -> str:
    mermaid_lines = ["graph TD"]
    
    for node in graph.nodes:
        if node.type == "concept":
            mermaid_lines.append(f"  {node.id}[{node.label}]:::concept")
        elif node.type == "definition":
            mermaid_lines.append(f"  {node.id}(\"{node.label}\"):::definition")
        elif node.type == "formula":
            mermaid_lines.append(f"  {node.id}[\"{node.label}\"]:::formula")
    
    for edge in graph.edges:
        mermaid_lines.append(f"  {edge.source} -- {edge.relationship} --> {edge.target}")
    
    mermaid_lines.extend([
        "  classDef concept fill:#4b5563,stroke:#1f2937,color:white",
        "  classDef definition fill:#6b7280,stroke:#374151,color:white",
        "  classDef formula fill:#9ca3af,stroke:#4b5563,color:white"
    ])
    
    return "\n".join(mermaid_lines)

def calculate_confidence(concept_data: Dict) -> float:
    num_definitions = len(concept_data.get("definitions", {}))
    if num_definitions == 0:
        return 0.3
    elif num_definitions == 1:
        return 0.7
    else:
        definitions = list(concept_data["definitions"].values())
        if all(d == definitions[0] for d in definitions):
            return 1.0
        else:
            return 0.5

# --- Written by ARUSH SHARMA ---